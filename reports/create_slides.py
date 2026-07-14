"""
Generate PowerPoint presentation for the Qualifyze case study.
Output: reports/presentation.pptx

Design: Clean, professional, minimal text per slide.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

REPORTS_DIR = Path(__file__).parent

# Design tokens
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x26, 0x6D, 0xF0)
TEXT_PRIMARY = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_SECONDARY = RGBColor(0x5A, 0x5A, 0x6E)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xCC, 0xCC, 0xDD)
SUCCESS = RGBColor(0x10, 0xB9, 0x81)
WARNING = RGBColor(0xF5, 0x9E, 0x0B)
DANGER = RGBColor(0xEF, 0x44, 0x44)


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=TEXT_PRIMARY, align=PP_ALIGN.LEFT, line_spacing=1.2):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = align
        p.space_after = Pt(size * 0.4)
    return tf


def title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_DARK)
    add_text(slide, Inches(1), Inches(2.2), Inches(8), Inches(1.5),
             title, size=38, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, Inches(1), Inches(3.8), Inches(8), Inches(1),
                 subtitle, size=16, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    return slide


def section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, ACCENT)
    add_text(slide, Inches(1), Inches(2.8), Inches(8), Inches(1.5),
             title, size=34, bold=True, color=TEXT_WHITE, align=PP_ALIGN.CENTER)
    return slide


def content_slide(prs, title, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_LIGHT)
    add_text(slide, Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.8),
             title, size=26, bold=True, color=TEXT_PRIMARY)

    # Subtle line under title
    shape = slide.shapes.add_shape(1, Inches(0.6), Inches(1.1), Inches(2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()

    bullet_text = "\n".join(f"  {b}" for b in bullets)
    add_text(slide, Inches(0.6), Inches(1.4), Inches(8.8), Inches(5),
             bullet_text, size=15, color=TEXT_PRIMARY)

    if note:
        add_text(slide, Inches(0.6), Inches(6.5), Inches(8.8), Inches(0.6),
                 note, size=11, color=TEXT_SECONDARY)
    return slide


def image_slide(prs, title, image_path, caption=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_LIGHT)
    add_text(slide, Inches(0.6), Inches(0.3), Inches(8.8), Inches(0.7),
             title, size=22, bold=True, color=TEXT_PRIMARY)

    if Path(image_path).exists():
        slide.shapes.add_picture(str(image_path), Inches(0.6), Inches(1.1), width=Inches(8.8))

    if caption:
        add_text(slide, Inches(0.6), Inches(6.8), Inches(8.8), Inches(0.5),
                 caption, size=10, color=TEXT_SECONDARY)
    return slide


def metrics_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_LIGHT)
    add_text(slide, Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.8),
             "Model Performance", size=26, bold=True, color=TEXT_PRIMARY)

    # Big metrics
    metrics = [
        ("ROC-AUC", "0.875", "Discrimination ability"),
        ("PR-AUC", "0.376", "7x better than random"),
        ("Recall@80%", "80%", "Catches 4/5 OAIs"),
    ]

    for i, (name, value, desc) in enumerate(metrics):
        x = Inches(0.8 + i * 3.2)
        add_text(slide, x, Inches(1.8), Inches(2.8), Inches(0.5),
                 name, size=13, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)
        add_text(slide, x, Inches(2.3), Inches(2.8), Inches(0.8),
                 value, size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_text(slide, x, Inches(3.3), Inches(2.8), Inches(0.5),
                 desc, size=11, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER)

    # Comparison table
    table_text = (
        "Model                         ROC-AUC   PR-AUC\n"
        "────────────────────────────────────────────────\n"
        "Logistic Regression            0.739     0.254\n"
        "XGBoost (weighted)             0.875     0.376  ← final\n"
        "XGBoost (tuned)                0.870     0.358\n"
    )
    add_text(slide, Inches(0.8), Inches(4.2), Inches(8.4), Inches(2.5),
             table_text, size=13, color=TEXT_PRIMARY)

    return slide


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 1 — Title
    slide = title_slide(prs, "Supplier Non-Compliance\nRisk Prediction",
                        "AI Data Scientist Case Study — Qualifyze\n\nSantiago Dominguez  |  July 2026")

    # 2 — Problem
    content_slide(prs, "The Problem", [
        "Qualifyze evaluates suppliers for pharma/device companies",
        "FDA inspects facilities → classifies as NAI (ok) / VAI (minor) / OAI (serious)",
        "OAI = enforcement actions, warning letters, recalls",
        "Only ~5% of inspections result in OAI — rare but very costly",
        "",
        "GOAL: Predict which suppliers will get OAI before it happens",
        "→ Enable proactive risk management instead of reactive response",
    ])

    # 3 — Data
    section_slide(prs, "Data & Processing")

    content_slide(prs, "Data Sources (all public FDA)", [
        "Inspection Classifications    337K records   NAI/VAI/OAI per facility",
        "Citations (483 observations)  277K records   Specific CFR violations",
        "Warning Letters               3,608 records  Formal enforcement notices",
        "Enforcement / Recalls         57K records    Product recalls by severity",
        "Published Form 483s           1,973 records  Detailed inspection reports",
        "",
        "All linked via FEI Number (facility ID)",
        "Warning letters + enforcement: fuzzy name matching (TF-IDF + location)",
    ], note="Coverage: WL 55% matched, Enforcement 82% matched. Remaining are online pharmacies with no inspection history.")

    # 5 — Features
    content_slide(prs, "Feature Engineering", [
        "Unit of analysis: one row per inspection",
        "All features from PRIOR history only (no data leakage)",
        "",
        "INSPECTION HISTORY    n_prior_oai, n_prior_nai, trend, recency",
        "CITATIONS             total count, unique CFR codes, max in single inspection",
        "WARNING LETTERS       count, days since last (via entity matching)",
        "ENFORCEMENT           recall count, Class I count",
        "483s                  count of published forms",
        "CONTEXT               product type, country, project area",
        "",
        "22 features total  |  Temporal split: train < 2025, test ≥ 2025",
    ])

    # 6 — Modeling approach
    section_slide(prs, "Modeling")

    content_slide(prs, "Approach & Decisions", [
        "Baseline: Logistic Regression (fully interpretable)",
        "Main: XGBoost with class weighting (scale_pos_weight = 20)",
        "Tuning: RandomizedSearchCV, 30 iterations, PR-AUC, TimeSeriesSplit",
        "",
        "KEY DECISIONS:",
        "  • No SMOTE — 7,200 positives is enough; class weights suffice",
        "  • Removed last_classification_oai — model improved without it",
        "    (was overfitting to a trivial pattern; other features encode same info)",
        "  • Threshold optimized for recall (asymmetric cost in supply chain risk)",
    ])

    # 7 — Metrics
    metrics_slide(prs)

    # 8 — ROC/PR
    image_slide(prs, "ROC & Precision-Recall Curves",
                REPORTS_DIR / "fig_model_comparison.png")

    # 9 — Threshold
    content_slide(prs, "Threshold: Business-Driven Choice", [
        "The model outputs a probability → we choose where to cut",
        "",
        "FALSE POSITIVE (flag safe supplier)  → analyst reviews it, low cost",
        "FALSE NEGATIVE (miss risky supplier) → client gets enforcement action, high cost",
        "",
        "Recall 80%  |  Precision 14%  |  ~6 false alarms per true alert",
        "",
        "This is a SCREENING tool — flags for human review, not auto-reject",
        "Threshold is configurable per client's review capacity",
    ])

    # 10 — Feature importance
    image_slide(prs, "What Drives Risk? (Feature Importance)",
                REPORTS_DIR / "fig_feature_importance.png",
                "n_prior_oai (41%) + recent_oai_rate (11%) + n_prior_nai (9%) — distributed, robust signal")

    # 11 — SHAP
    section_slide(prs, "Interpretability")

    image_slide(prs, "SHAP — Every Prediction is Explainable",
                REPORTS_DIR / "fig_shap_beeswarm.png",
                "Each dot = one inspection. Color = feature value (red=high). Position = impact on risk.")

    # 12 — Risk report example
    content_slide(prs, "Supplier Risk Report (Production Output)", [
        "FACILITY: Acme Pharma (FEI: 3007058211)",
        "RISK SCORE: 95.2%  —  CRITICAL",
        "",
        "Risk Drivers:",
        "  ↑ 2 prior OAI inspections              (+25%)",
        "  ↑ Recent OAI rate = 1.0                 (+15%)",
        "  ↑ 0 clean (NAI) inspections             (+10%)",
        "  ↑ Veterinary / high-risk sector          (+5%)",
        "",
        "Every prediction has a human-readable explanation",
        "Full audit trail for regulatory compliance",
    ])

    # 13 — Deployment
    section_slide(prs, "Production & MLOps")

    content_slide(prs, "Deployed System", [
        "LIVE APP: Streamlit on Azure App Service",
        "  → Enter FEI → get risk score + SHAP explanation",
        "  → qualifyze-risk-xxx.azurewebsites.net",
        "",
        "CI/CD: GitHub Actions",
        "  → PR → pytest → merge → tests pass → auto-deploy to Azure",
        "  → Branch protection on main",
        "",
        "ORCHESTRATION: Airflow (4 DAGs, local demo)",
        "  → Weekly data refresh, feature pipeline, monthly retrain, batch scoring",
        "  → Champion/challenger model promotion pattern",
        "",
        "MONITORING: drift detection on score distribution (proposed)",
    ])

    # 14 — Limitations
    content_slide(prs, "Limitations & Next Steps", [
        "LIMITATIONS:",
        "  • Cannot predict first-time offenders (no prior history)",
        "  • FDA inspection database is not comprehensive",
        "  • Entity matching has noise for ~45% of warning letters",
        "  • ROC-AUC 0.87 is likely the ceiling with public data alone",
        "",
        "HIGHEST-VALUE NEXT STEPS:",
        "  1. Integrate Qualifyze private audit data (biggest lift for first-offenders)",
        "  2. NLP on citation text (severity categorization beyond CFR codes)",
        "  3. Survival analysis (time-to-next-OAI, not just binary)",
        "  4. Active learning from analyst feedback on flagged suppliers",
    ])

    # 15 — Closing
    slide = title_slide(prs, "Thank You",
                        "Questions & Discussion\n\ngithub.com/santiago-dc/qualifyze-case-study")

    # Save
    output = REPORTS_DIR / "presentation.pptx"
    prs.save(str(output))
    print(f"Saved: {output} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    create_presentation()
